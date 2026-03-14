import os
import logging

logging.disable(logging.DEBUG)  # 关闭DEBUG日志的打印
logging.disable(logging.WARNING)  # 关闭WARNING日志的打印

import subprocess
import tkinter as tk
import pptx
from pptx import Presentation
from docx import Document
from openpyxl import load_workbook
import openpyxl
from openpyxl.styles import PatternFill, Font, Alignment
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
import os
import win32com.client as win32
from win32com.client import constants
from ocr import *
from PIL import ImageGrab

def run_powershell(command):
    """
    新建一个终端然后运行powershell命令,你可以以管理员身份执行大部分命令包括操作读写文件,安装软件等等,注意加以限制以避免输出过多信息

    Args:
      command (str):powershell命令

    Returns:
      powershell输出 (str)
    """
    try:
        process = subprocess.Popen(
            ["powershell", "-Command", command],
            shell=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
        )
        output, errors = process.communicate()

        opt = output.decode("GBK").replace("  ", " ")
        for i in range(5):
            opt = opt.replace("  ", " ")

        res = (
            "\n输出:\n"
            + (
                opt
                if len(opt) < 2500
                else opt[:2500] + "...(信息过多,无法完全显示)"
            )
            + ("无" if len(opt) == 0 else "")
            + (
                "\n状态:运行成功"
                if errors.decode("GBK") == ""
                else "\n错误:\n" + errors.decode("GBK")
            )
        )
        return res
    except UnicodeDecodeError:
        try:
            process = subprocess.Popen(
                ["powershell", "-Command", command],
                shell=True,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
            )
            output, errors = process.communicate()

            opt = output.decode("utf-8").replace("  ", " ")
            for i in range(5):
                opt = opt.replace("  ", " ")

            res = (
                "\n输出:\n"
                + (opt if len(opt) < 2500 else opt[:2500] + "...(信息过多,无法完全显示)")
                + ("无" if len(opt) == 0 else "")
                + (
                    "\n状态:运行成功"
                    if errors.decode("utf-8") == ""
                    else "\n错误:\n" + errors.decode("utf-8")
                )
            )
            return res
        except Exception as e:
            return f"错误: {str(e)}"
    except Exception as e:
        return f"错误: {str(e)}"

def wait_user_do(operation):
    """
    请求并等待用户操作完成(此工具用于等待用户操作,不能用于询问用户)

    Args:
      operation (str):操作的方法,比如"请扫描二维码"

    Returns:
      操作结果
    """

    window = tk.Tk()
    window.attributes("-topmost", True)
    window.attributes("-alpha", 0.6)
    window.title("完成后请关闭窗口")
    w, h = window.winfo_screenwidth(), window.winfo_screenheight()
    window.geometry(
        "%dx%d+%d+%d" % (w // 4, h // 8, w // 2 - w // 8, h // 2.5 - h // 16)
    )
    # window.withdraw()
    show_text = tk.Label(window, font=("微软雅黑", 13), text="", anchor="center")
    show_text.place(relx=0.5, rely=0.5, anchor="center")

    operation2 = operation
    if "\n" not in operation:
        lines = []
        for i in range(0, len(operation), 20):
            lines.append(operation[i : i + 20])
        operation2 = "\n".join(lines)
    show_text.config(text=operation2)

    window.mainloop()

    return "用户操作完成"


def read_file(file_path, encoding="utf-8"):
    """
    读取文件内容

    参数:
        file_path (str): 文件路径
        encoding (str): 文件编码格式，默认为utf-8

    返回:
        str: 文件内容或错误信息
    """
    try:
        with open(file_path, "r", encoding=encoding) as file:
            content = file.read()

            if len(content) > 3000:
                content = content[:3000] + "\n...(信息过多,无法完全显示)"

            return content
    except FileNotFoundError:
        return f"错误: 文件 '{file_path}' 未找到"
    except PermissionError:
        return f"错误: 没有权限访问文件 '{file_path}'"
    except Exception as e:
        return f"错误: 读取文件时发生异常 - {str(e)}"


def create_file(text, path):
    """
    将text写入path对应的文件(utf-8编码)

    Args:
      text (str):文件内容
      name (str):文件名,比如"D:/a.py","D:/b.txt","D:/c.md","D:/d.html"等等

    Returns:
      status (str):执行状态
    """
    with open(path, "w", encoding="utf-8") as f:
        f.write(text)

    return "%s创建成功" % (path)


def run_python_code(code, name="temp.py"):
    """
    临时保存并运行python代码来完成各种任务包括操作excel,word,ppt,批量处理本地文件,绘制图表,开发程序等等,如果缺失模块可以自行安装,文件会自动覆盖保存到"D:\\ExternalFiles\\{name}",你可以通过run_powershell工具多次重新运行或者复制到其他位置
    注意: 要在程序里面print一些状态比如"xxx操作已完成","xxx错误"等等来方便检查当前代码的执行情况,否则输出将为空
    注意: 程序中使用的所有路径必须为绝对路径
    **强烈建议: 尽可能用create_python_tool把程序制作成CLI工具而不是临时文件,使其能够解决一类问题,这样可以方便以后重复使用,提高效率**
    此外,**如果出现错误,可以使用replace_file_content工具快速修复再`cd D:/ExternalFiles/;python {name}.py`重新运行**

    Args:
      code (str):python代码内容
      name (str):文件名,比如"a.py"

    Returns:
      opt (str):程序输出,不报错即成功
    """
    with open("D:/ExternalFiles/%s" % (name), "w", encoding="utf-8") as f:
        f.write(code)

    return run_powershell("python D:/ExternalFiles/%s" % (name))

def create_python_tool(tool_name, code, instructions):
    """
    创建外部python文件作为CLI工具(推荐)
    python文件会自动覆盖保存到"D:\\ExternalFiles\\{tool_name.py}"
    CLI工具的使用说明会自动覆盖保存到"D:\\ExternalFiles\\{tool_name.md}"
    创建完之后不会自动运行,**要通过调用run_powershell工具执行命令`cd D:/ExternalFiles/;python {name}.py 具体启动参数`来运行**
    注意: 要在程序里面print一些状态比如"xxx操作已完成","xxx错误"等等来方便检查当前工具的执行情况,否则输出将为空

    Args:
      tool_name (str):工具名,不是文件名,比如"TOOL_xxx",尽量准确且达意
      code (str):python代码内容,需要读取CLI启动参数作为程序输入
      instructions (str):CLI工具的详细使用说明,markdown格式

    Returns:
      opt (str):程序输出,不报错即成功
    """
    with open("D:/ExternalFiles/%s.py" % (tool_name), "w", encoding="utf-8") as f:
        f.write(code)
    with open("D:/ExternalFiles/%s.md" % (tool_name), "w", encoding="utf-8") as f:
        f.write(instructions)

    return f"工具创建成功,可以通过调用run_powershell工具执行命令`cd D:/ExternalFiles/;python {tool_name}.py 具体启动参数`来运行"

def replace_file_content(file_path, old_text, new_text, backup=False):
    """
    替换指定文件中的部分内容,支持文本模式的文件包括txt,csv,json,xml,py等,**推荐使用此工具快速修复python代码的bug**

    参数:
        file_path (str): 文件路径
        old_text (str): 需要被替换的文本
        new_text (str): 替换后的文本
        backup (bool): 是否创建备份文件，默认为False

    返回:
        str: 替换状态
    """
    try:
        if not os.path.exists(file_path):
            return "文件不存在"

        with open(file_path, "r", encoding="utf-8") as file:
            content = file.read()

        if old_text not in content:
            return f"警告: 文本 '{old_text}' 未在文件中找到"

        if backup:
            backup_path = file_path + ".bak"
            with open(backup_path, "w", encoding="utf-8") as backup_file:
                backup_file.write(content)

        new_content = content.replace(old_text, new_text)

        with open(file_path, "w", encoding="utf-8") as file:
            file.write(new_content)

        return True

    except FileNotFoundError as e:
        return f"错误: {e}"
    except PermissionError as e:
        return f"权限错误: {e}"
    except Exception as e:
        return f"替换过程中发生错误: {e}"


def replace_ppt_content(old_file_path, new_file_path, replace_list):
    """
    替换PPT文件中的文本内容,可用于修改ppt模板文字,ppt润色,在此之前,你应该先通过read_ppt_and_export_txt工具读取ppt内容

    参数:
    old_file_path: 旧PPT文件路径
    new_file_path: 新PPT文件路径
    replace_list: 替换列表，格式为[("旧文本1","新文本1"),("旧文本2","新文本2")...]

    返回:
    str: 状态
    """
    try:
        # 打开旧的PPT文件
        prs = Presentation(old_file_path)

        # 遍历所有幻灯片
        for slide in prs.slides:
            # 遍历幻灯片中的所有形状
            for shape in slide.shapes:
                # 检查形状是否包含文本
                if hasattr(shape, "text") and shape.text:
                    # 对每个替换项进行文本替换
                    for old_text, new_text in replace_list:
                        if old_text in shape.text:
                            # 替换文本
                            shape.text = shape.text.replace(old_text, new_text)

                # 处理文本框中的文本
                if hasattr(shape, "text_frame") and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text:
                                for old_text, new_text in replace_list:
                                    if old_text in run.text:
                                        run.text = run.text.replace(old_text, new_text)

        # 保存新的PPT文件
        prs.save(new_file_path)
        return f"成功保存新文件: {new_file_path}"

    except Exception as e:
        return f"处理PPT文件时出错: {str(e)}"


def create_ppt_from_txt(path, output_file):
    """
    根据读取特定格式的文本文件创建 PowerPoint 演示文稿。

    参数:
        path (str): 文本文件的路径。
        具体格式:
        1.使用 --- 作为幻灯片之间的分隔符
        2.每张幻灯片的第一行非空行会被作为标题,可以以#开头
        3.标题之后的所有行视为幻灯片内容,每行对应一个段落
        4.以 - 或 * 开头的行（注意要有后面的空格）会被渲染为带项目符号的列表
        5.空行可用于增加段落间的垂直间距
        output_file (str): 输出的 .pptx 文件路径。
        注意每一页ppt连标题最好不要超过8行,否则可能无法完整显示

    返回:
        str: 状态
    """
    prs = Presentation()
    # 使用“标题和内容”布局（通常索引为 1）
    slide_layout = prs.slide_layouts[1]

    text = read_file(path)

    # 按分隔符 '---' 分割幻灯片
    slides_raw = []
    current_lines = []
    for line in text.splitlines():
        if line.strip() == "---":
            if current_lines:
                slides_raw.append(current_lines)
                current_lines = []
        else:
            current_lines.append(line)
    if current_lines:
        slides_raw.append(current_lines)

    for slide_lines in slides_raw:
        # 找到第一个非空行作为标题
        title = None
        content_lines = []
        idx = 0
        while idx < len(slide_lines) and slide_lines[idx].strip() == "":
            idx += 1
        if idx < len(slide_lines):
            title_line = slide_lines[idx].strip()
            if title_line.startswith("# ") or title_line.startswith("## "):
                title = title_line[2:].strip()
            else:
                title = title_line
            content_lines = slide_lines[idx + 1 :]
        else:
            continue  # 幻灯片无有效内容，跳过

        # 添加幻灯片
        slide = prs.slides.add_slide(slide_layout)

        # 设置标题
        if title:
            slide.shapes.title.text = title

        # 获取正文占位符（通常为第二个占位符）
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()  # 清除默认占位文本

        # 设置默认段落格式
        tf.paragraphs[0].text = ""  # 清空第一个段落

        # 解析内容行
        for i, line in enumerate(content_lines):
            stripped = line.strip()

            # 处理空行：添加空段落
            if stripped == "":
                p = tf.add_paragraph()
                p.text = ""
                continue

            # 检查是否为列表项（去除左边空格后判断）
            lstripped = line.lstrip()
            is_bullet = lstripped.startswith("- ") or lstripped.startswith("* ")

            # 添加新段落
            p = tf.add_paragraph()
            p.text = lstripped[2:].strip() if is_bullet else stripped

            if is_bullet:
                # 列表项：带项目符号
                p.bullet = True
                p.level = 0
            else:
                # 普通段落：不带项目符号
                p.bullet = False

                # **使用更安全的方法移除项目符号**
                try:
                    # 确保 pPr 元素存在
                    if not hasattr(p._element, "pPr") or p._element.pPr is None:
                        # 创建 pPr 元素
                        pPr = p._element.makeelement("pPr")
                        p._element.append(pPr)

                    # 移除现有的项目符号相关元素
                    pPr = p._element.pPr
                    for child in list(pPr):
                        if (
                            child.tag.endswith("buFont")
                            or child.tag.endswith("buChar")
                            or child.tag.endswith("buNone")
                        ):
                            pPr.remove(child)

                    # 添加 buNone 元素
                    buNone = pPr.makeelement("buNone")
                    pPr.append(buNone)

                except Exception as e:
                    # print(f"处理段落时出现警告: {e}")
                    # 备用方法：使用 XML 字符串
                    try:
                        pPr = p._element.get_or_add_pPr()
                        buNone = parse_xml(r"<a:buNone %s/>" % nsdecls("a"))
                        pPr.append(buNone)
                    except:
                        pass  # 如果都失败，至少保持了基本功能

    prs.save(output_file)
    return "创建成功"


def read_ppt_and_export_txt(ppt_path, txt_path=None):
    """
    读取PPT文件并导出所有内容到TXT文本文件

    参数:
        ppt_path (str): PPT文件路径
        txt_path (str, optional): 输出TXT文件路径，默认为PPT文件名 + .txt

    返回:
        str: 执行状态信息
    """

    try:
        # 检查PPT文件是否存在
        if not os.path.exists(ppt_path):
            return f"错误: PPT文件不存在 - {ppt_path}"

        # 默认输出文件名
        if txt_path is None:
            base_name = os.path.splitext(ppt_path)[0]
            txt_path = f"{base_name}.txt"

        # 读取PPT文件
        presentation = pptx.Presentation(ppt_path)

        # 准备输出内容
        content = []
        content.append(f"PPT文件内容提取报告")
        content.append(f"文件路径: {ppt_path}")
        content.append(f"总页数: {len(presentation.slides)}")
        content.append("=" * 50)
        content.append("")

        # 遍历每一页幻灯片
        for slide_index, slide in enumerate(presentation.slides, 1):
            content.append(f"第 {slide_index} 页幻灯片")
            content.append("-" * 30)

            # 获取幻灯片标题
            title = slide.shapes.title.text if slide.shapes.title else "无标题"
            content.append(f"标题: {title}")

            # 获取所有文本内容
            text_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content.append(shape.text.strip())

            # 添加文本内容
            if text_content:
                content.append("文本内容:")
                for i, text in enumerate(text_content, 1):
                    content.append(f"  {i}. {text}")
            else:
                content.append("无文本内容")

            content.append("")

        # 写入TXT文件
        with open(txt_path, "w", encoding="utf-8") as f:
            f.write("\n".join(content))
        for i in content:
            print(i)
        return f"成功: PPT内容已导出到 {txt_path}\n总共导出 {len(presentation.slides)} 页幻灯片内容"

    except ImportError:
        return "错误: 请安装python-pptx库 - pip install python-pptx"
    except Exception as e:
        return f"错误: {str(e)}"


def read_word_and_export_txt(word_path, txt_path=None):
    """
    读取Word文件并导出为txt文本

    参数:
        word_path (str): Word文件路径
        txt_path (str, optional): 输出txt文件路径，如果为None则自动生成

    返回:
        str: 执行状态信息
    """
    try:
        # 检查文件是否存在
        if not os.path.exists(word_path):
            return f"错误：文件 '{word_path}' 不存在"

        # 检查是否为Word文件
        if not word_path.lower().endswith((".doc", ".docx")):
            return "错误：请选择Word文件(.doc或.docx)"

        # 读取Word文件
        doc = Document(word_path)

        # 获取所有段落内容
        content_list = []
        for i, para in enumerate(doc.paragraphs):
            if para.text.strip():  # 只处理非空段落
                content_list.append(f"段落 {i+1}: {para.text}")

        # 获取所有表格内容
        table_content = []
        for i, table in enumerate(doc.tables):
            table_content.append(f"\n表格 {i+1}:")
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    row_text.append(cell.text.strip())
                table_content.append("  ".join(row_text))

        # 组合所有内容
        all_content = content_list + table_content

        # 如果没有指定输出路径，自动生成
        if txt_path is None:
            name, _ = os.path.splitext(word_path)
            txt_path = f"{name}_output.txt"

        # 写入txt文件
        with open(txt_path, "w", encoding="utf-8") as f:
            f.write("Word文件内容导出\n")
            f.write("=" * 50 + "\n")
            f.write(f"源文件: {os.path.basename(word_path)}\n")
            f.write("=" * 50 + "\n\n")

            for content in all_content:
                f.write(content + "\n")

        # 打印详细内容
        print("=" * 50)
        print("Word文件内容详情:")
        print("=" * 50)
        for content in all_content:
            print(content)
        print("=" * 50)
        print(f"文件已成功导出到: {txt_path}")
        print("=" * 50)

        return f"成功：文件已导出到 {txt_path}"

    except Exception as e:
        return f"错误：处理文件时出现异常 - {str(e)}"


def read_excel_and_export_txt(excel_file, txt_file=None):
    """
    读取Excel文件并将其内容导出为txt文本,**注意第一行是工作表名比如sheet1**

    参数:
        excel_file (str): Excel文件路径
        txt_file (str, optional): 输出txt文件路径，如果为None则使用excel_file同目录下的同名txt文件

    返回:
        str: 执行状态信息
    """
    try:
        # 加载Excel文件
        wb = load_workbook(excel_file)

        # 如果没有指定txt文件路径，使用excel文件同目录下的同名txt文件
        if txt_file is None:
            import os

            txt_file = os.path.splitext(excel_file)[0] + ".txt"

        # 打开txt文件准备写入
        with open(txt_file, "w", encoding="utf-8") as f:
            # 遍历所有工作表
            for sheet_name in wb.sheetnames:
                f.write(f"=== 工作表: {sheet_name} ===\n")

                # 获取工作表
                ws = wb[sheet_name]

                # 读取并写入所有数据
                for row in ws.iter_rows(values_only=True):
                    row_data = []
                    for cell in row:
                        if cell is None:
                            row_data.append("")
                        else:
                            row_data.append(str(cell))
                    f.write("\t".join(row_data) + "\n")

                f.write("\n")  # 工作表之间添加空行

            # f.write(f"文件已成功导出到: {txt_file}")

        return f"Excel文件读取成功，数据已导出到: {txt_file}"

    except FileNotFoundError:
        return f"错误: 找不到文件 {excel_file}"
    except Exception as e:
        return f"错误: {str(e)}"


def create_excel_from_2d_list(data, filename="output.xlsx", sheet_name="Sheet1"):
    """
    通过二维列表创建Excel文件

    参数:
    data: 二维列表，每个元素是元组(value, 样式, 背景色)或数字,样式可以为""或"bold"或"italic"或"underline";背景色可以为""(无色)或十六进制颜色代码
    filename: 输出文件名，默认为"output.xlsx"
    sheet_name: 工作表名称，默认为"Sheet1"

    返回:
    字符串，表示执行状态
    """
    try:
        # 创建工作簿和工作表
        workbook = openpyxl.Workbook()
        worksheet = workbook.active
        worksheet.title = sheet_name

        # 遍历二维列表
        for row_idx, row_data in enumerate(data, 1):
            for col_idx, cell_data in enumerate(row_data, 1):
                # 处理不同的数据类型
                if isinstance(cell_data, tuple) and len(cell_data) >= 1:
                    value = cell_data[0]
                    style = cell_data[1] if len(cell_data) > 1 else None
                    bg_color = cell_data[2] if len(cell_data) > 2 else None
                else:
                    # 处理数字或字符串
                    value = cell_data
                    style = None
                    bg_color = None

                # 写入单元格
                cell = worksheet.cell(row=row_idx, column=col_idx, value=value)

                # 应用样式
                if style:
                    if "bold" in style:
                        cell.font = Font(name="等线",bold=True)
                    if "italic" in style:
                        cell.font = Font(name="等线", italic=True)
                    if "underline" in style:
                        cell.font = Font(name="等线", underline="single")
                else:
                    cell.font = Font(name="等线")

                # 应用背景色
                if bg_color:
                    try:
                        fill = PatternFill(
                            start_color=bg_color.lstrip("#"),
                            end_color=bg_color.lstrip("#"),
                            fill_type="solid",
                        )
                        cell.fill = fill
                    except:
                        # 如果颜色格式错误，忽略背景色设置
                        pass

                # 自动调整列宽
                worksheet.column_dimensions[
                    openpyxl.utils.get_column_letter(col_idx)
                ].width = 15

        # 保存文件
        workbook.save(filename)
        return f"成功创建Excel文件: {filename}"

    except Exception as e:
        return f"创建Excel文件失败: {str(e)}"


def convert_word_or_txt_to_pdf(word_path, pdf_path):
    """
    将word或txt文件转换为PDF格式,注意路径要用反斜杠

    参数:
        word_path (str): 输入文件的路径
        pdf_path (str, optional): 输出PDF文件的路径

    返回:
        str: 转换状态信息
    """
    # 创建Word应用程序实例
    word = win32.gencache.EnsureDispatch("Word.Application")

    try:
        doc = word.Documents.Open(word_path)
        doc.SaveAs(pdf_path, FileFormat=17)  # 17是PDF格式的文件格式代码
        doc.Close()
        word.Quit()

        return "转换成功"
    except Exception as e:
        # 退出Word应用程序
        word.Quit()
        return str(e)


def convert_ppt_to_pdf(ppt_path, pdf_path):
    """
    将ppt文件转换为PDF格式,注意路径要用反斜杠

    参数:
        ppt_path (str): 输入文件的路径
        pdf_path (str, optional): 输出PDF文件的路径

    返回:
        str: 转换状态信息
    """
    # 创建Word应用程序实例
    ppt = win32.gencache.EnsureDispatch("PowerPoint.Application")

    try:
        doc = ppt.Presentations.Open(ppt_path)
        doc.SaveAs(pdf_path, FileFormat=32)
        doc.Close()
        ppt.Quit()

        return "转换成功"
    except Exception as e:
        # 退出Word应用程序
        ppt.Quit()
        return str(e)


def add_knowledge(knowledge):
    """
    记录此次回答获取的知识,包括用户的习惯,系统的信息(比如电脑用户名,桌面路径,硬件参数等),解决问题的技巧等任何对之后的回答有利的信息

    参数:
        knowledge (str): 一句话概括此次回答获取的知识,**要求简短且有效**

    返回:
        str: 执行状态信息
    """
    try:
        with open("D:/ExternalFiles/KNOWLEDGE.txt", "a", encoding="utf-8") as file:
            file.write(knowledge + "\n")
        return f"成功添加知识"
    except Exception as e:
        return f"添加内容失败: {str(e)}"


def get_screen_image(save_path):
    """
    获取当前电脑屏幕图像并保存,**不要乱用此工具!,只有需要提取屏幕上的文字时才使用**

    Args:
        save_path (str): 保存路径

    Returns:
        str: 图像保存路径或错误信息
    """
    try:
        screenshot = ImageGrab.grab().resize((1920, 1200))
        screenshot.save(save_path)  # 保存截图
        return "截图成功"
    except ImportError:
        return "错误: 未安装 PIL 库，请运行 'pip install pillow'"
    except Exception as e:
        return f"错误: {str(e)}"


if __name__ == "__main__":
    # create_file("＜（＾－＾）＞,🐕","./a.md")
    # result = convert_word_to_pdf(
    #     "D:\Python_items\LocalAgent\output.docx", "./output.pdf"
    # )
    # result = convert_word_or_txt_to_pdf(
    #     "D:\\Python_items\\LocalAgent\\a.txt", "D:\\Python_items\\LocalAgent\\a.pdf"
    # )
    # print(result)
    # sample_data = [
    #     [
    #         ("姓名", "bold", "FFCCCB"),
    #         ("年龄", "bold", "FFCCCB"),
    #         ("城市", "bold", "FFCCCB"),
    #     ],
    #     [("张三", "italic", "FFFFFF"), (25, None, "FFFFFF"), ("北京", None, "FFFFFF")],
    #     [
    #         ("李四", "underline", "FFFFFF"),
    #         (30, None, ""),
    #         ("上海", None, ""),
    #     ],
    #     [100, 200, 300],  # 纯数字
    #     [("测试", "bold", "FFFF00"), ("数据", "italic", "FFFF00"), 456],
    # ]
    # # 调用函数
    # result = create_excel_from_2d_list(sample_data, "test.xlsx")
    # print(result)
    # replace_ppt_content(
    #     "C:\\Users\\22974\\Desktop\\热门.pptx",
    #     "C:\\Users\\22974\\Desktop\\热门2.pptx",
    #     [("我们","你们")]
    # )
    # print(run_powershell("ls && ls"))
    # read_excel_and_export_txt("C:/Users/22974/Desktop/data.xlsx")
    # result = convert_markdown_to_word("test.md", "output.docx")
    # print(result)
    # print(read_word_and_export_txt("C:\\Users\\22974\\Downloads\\参考答案（教师版）.docx"))
    # print(read_ppt_and_export_text("C:\\Users\\22974\\Desktop\\热门.pptx"))
    # print(run_powershell("cd D:/python_items;ls"))
    # press_key("win + s")
    # 生成 PPT
    # create_ppt_from_txt("./a.txt", "demo.pptx")
    # print(run_python_code("print(45616)"))
    # replace_file_content("D:\\Temp\\txt_to_pdf.py", "TXT文件","TEXT文件")
    # print(read_ppt_and_export_txt("C:\\Users\\22974\\Downloads\\平衡状态的判断.ppt","res.txt"))
    # print()
    # print(create_python_tool("echo","import sys;print('参数',sys.argv[0])","# 打印"))
    # add_knowledge("桌面路径: C:/Users/22974/Desktop")
    # print(get_screen_image("./a.jpg"))
    pass
